"""Generate styled PowerPoint decks from the curriculum Markdown sources."""

from __future__ import annotations

import argparse
import re
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DECKS_DIR = REPO_ROOT / "decks"
OUTLINE_PATH = DECKS_DIR / "token-optimization-context-engineering.outline.md"
NOTES_PATH = DECKS_DIR / "token-optimization-context-engineering.speaker-notes.md"

SLIDE_W = 13.333
SLIDE_H = 7.5


@dataclass(frozen=True)
class SlideSpec:
    number: int
    title: str
    bullets: list[str]
    notes: str


@dataclass(frozen=True)
class Theme:
    name: str
    label: str
    output_name: str
    background: RGBColor
    card: RGBColor
    card_alt: RGBColor
    title: RGBColor
    body: RGBColor
    muted: RGBColor
    accent: RGBColor
    accent2: RGBColor
    inverse: bool = False
    mono: bool = False


def color(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


THEMES: dict[str, Theme] = {
    "workshop": Theme(
        name="workshop",
        label="Workshop",
        output_name="token-optimization-context-engineering.pptx",
        background=color("#f6f8fa"),
        card=color("#ffffff"),
        card_alt=color("#eef4ff"),
        title=color("#24292f"),
        body=color("#57606a"),
        muted=color("#6e7781"),
        accent=color("#0969da"),
        accent2=color("#8250df"),
    ),
    "executive": Theme(
        name="executive",
        label="Executive briefing",
        output_name="token-optimization-context-engineering.executive.pptx",
        background=color("#0d1117"),
        card=color("#161b22"),
        card_alt=color("#1f2937"),
        title=color("#f0f6fc"),
        body=color("#c9d1d9"),
        muted=color("#8b949e"),
        accent=color("#58a6ff"),
        accent2=color("#a371f7"),
        inverse=True,
    ),
    "technical": Theme(
        name="technical",
        label="Technical deep dive",
        output_name="token-optimization-context-engineering.technical.pptx",
        background=color("#0b1020"),
        card=color("#111827"),
        card_alt=color("#172033"),
        title=color("#e6edf3"),
        body=color("#b7c3d0"),
        muted=color("#7d8590"),
        accent=color("#3fb950"),
        accent2=color("#ffab70"),
        inverse=True,
        mono=True,
    ),
}


def parse_outline(path: Path) -> list[tuple[int, str, list[str]]]:
    slide_re = re.compile(r"^## Slide (\d+):\s*(.+?)\s*$")
    slides: list[tuple[int, str, list[str]]] = []
    current_number: int | None = None
    current_title = ""
    current_bullets: list[str] = []

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        match = slide_re.match(line)
        if match:
            if current_number is not None:
                slides.append((current_number, current_title, current_bullets))
            current_number = int(match.group(1))
            current_title = match.group(2)
            current_bullets = []
            continue
        if current_number is None:
            continue
        if line.startswith("- "):
            current_bullets.append(line[2:].strip())

    if current_number is not None:
        slides.append((current_number, current_title, current_bullets))

    if not slides:
        raise ValueError(f"No slides found in {path}")
    return slides


def parse_notes(path: Path) -> dict[int, str]:
    note_re = re.compile(r"^## Slide (\d+)\s*$")
    notes: dict[int, list[str]] = {}
    current_number: int | None = None

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.rstrip()
        match = note_re.match(line.strip())
        if match:
            current_number = int(match.group(1))
            notes[current_number] = []
            continue
        if current_number is not None:
            notes[current_number].append(line)

    parsed = {number: "\n".join(lines).strip() for number, lines in notes.items()}
    if not parsed:
        raise ValueError(f"No speaker notes found in {path}")
    return parsed


def load_slides() -> list[SlideSpec]:
    outline = parse_outline(OUTLINE_PATH)
    notes = parse_notes(NOTES_PATH)
    specs: list[SlideSpec] = []

    for number, heading, bullets in outline:
        title = bullets[0] if heading == "Title" and bullets else heading
        body = bullets[1:] if heading == "Title" and bullets else bullets
        note = notes.get(number, "")
        if not note:
            raise ValueError(f"Missing speaker notes for slide {number}: {title}")
        specs.append(SlideSpec(number=number, title=title, bullets=body, notes=note))

    expected = list(range(1, len(specs) + 1))
    actual = [slide.number for slide in specs]
    if actual != expected:
        raise ValueError(f"Slides must be sequential. Expected {expected}, found {actual}")
    return specs


def font_name(theme: Theme, display: bool = False) -> str:
    if theme.mono and not display:
        return "Cascadia Code"
    return "Aptos Display" if display else "Aptos"


def set_fill(shape, fill_color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color


def add_shape(slide, kind, x: float, y: float, w: float, h: float, fill: RGBColor):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    theme: Theme,
    *,
    size: int = 20,
    bold: bool = False,
    color_value: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
    display: bool = False,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name(theme, display)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color_value or theme.body
    p.alignment = align
    return shape


def add_multiline_text(
    shape,
    lines: list[str],
    theme: Theme,
    *,
    size: int = 18,
    color_value: RGBColor | None = None,
    numbered: bool = False,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"{index + 1}. {line}" if numbered else line
        p.font.name = font_name(theme)
        p.font.size = Pt(size)
        p.font.color.rgb = color_value or theme.body
        p.space_after = Pt(7)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    theme: Theme,
    *,
    fill: RGBColor | None = None,
    title: str | None = None,
    body: list[str] | None = None,
    title_size: int = 15,
    body_size: int = 16,
    numbered: bool = False,
):
    card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill or theme.card)
    if title:
        add_text(slide, title, x + 0.22, y + 0.17, w - 0.44, 0.36, theme, size=title_size, bold=True, color_value=theme.title)
        if body:
            body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.62), Inches(w - 0.36), Inches(h - 0.74))
            add_multiline_text(body_box, body, theme, size=body_size, numbered=numbered)
    elif body:
        add_multiline_text(card, body, theme, size=body_size, numbered=numbered)
    return card


def apply_background(slide, theme: Theme) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme.background


def add_header(slide, spec: SlideSpec, theme: Theme) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.16, theme.accent)
    add_text(slide, spec.title, 0.65, 0.38, 10.7, 0.55, theme, size=25, bold=True, color_value=theme.title, display=True)
    add_text(slide, theme.label, 11.0, 0.45, 1.7, 0.3, theme, size=9, bold=True, color_value=theme.accent, align=PP_ALIGN.RIGHT)


def add_footer(slide, index: int, total: int, theme: Theme) -> None:
    add_text(
        slide,
        f"Token Optimization + Context Engineering | {index} / {total}",
        0.65,
        6.96,
        7.2,
        0.28,
        theme,
        size=9,
        color_value=theme.muted,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.35, 7.08, 4.3, 0.05, theme.card_alt)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.35, 7.08, 4.3 * (index / total), 0.05, theme.accent)


def add_notes(slide, spec: SlideSpec) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    paragraphs = [f"Slide {spec.number}: {spec.title}"]
    paragraphs.extend([part.strip() for part in re.split(r"\n\s*\n", spec.notes) if part.strip()])
    for index, paragraph in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = paragraph
        p.font.name = "Aptos"
        p.font.size = Pt(12)


def add_title_slide(slide, spec: SlideSpec, theme: Theme, total: int) -> None:
    apply_background(slide, theme)
    if not theme.inverse:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, color("#ffffff"))
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 4.25, SLIDE_H, theme.accent)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 4.25, 0, 0.18, SLIDE_H, theme.accent2)
        add_text(slide, "Token\nOptimization", 0.62, 1.25, 2.9, 1.35, theme, size=26, bold=True, color_value=color("#ffffff"), display=True)
        add_text(slide, "Context engineering turns Copilot usage into a repeatable design practice.", 0.68, 4.95, 2.95, 0.9, theme, size=16, color_value=color("#eaf2ff"))
        title_x = 4.95
        title_w = 7.25
        subtitle_x = 4.98
        subtitle_w = 6.95
        title_color = theme.title
        body_color = theme.body
    else:
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, theme.background)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 1.25, theme.card)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.25, SLIDE_W, 0.12, theme.accent)
        title_x = 0.72
        title_w = 7.0
        subtitle_x = 0.78
        subtitle_w = 4.95
        title_color = theme.title
        body_color = theme.body

    add_text(slide, "Workshop", title_x, 0.78, 2.2, 0.32, theme, size=12, bold=True, color_value=theme.accent2)
    add_text(slide, spec.title, title_x, 1.6, title_w, 1.55, theme, size=38, bold=True, color_value=title_color, display=True)
    add_text(slide, "\n".join(spec.bullets), subtitle_x, 3.42, subtitle_w, 1.0, theme, size=18, color_value=body_color)

    labels = ["Context hygiene", "Prompt discipline", "Model/surface routing", "Tool scope", "Measurement"]
    for index, label in enumerate(labels):
        y = 1.28 + index * 0.9
        fill = theme.card if theme.inverse else color("#f6f8fa")
        add_card(slide, 8.0, y, 4.4, 0.58, theme, fill=fill, title=f"{index + 1}. {label}", body=[], title_size=14)

    add_text(slide, "Generated from Markdown source with embedded speaker notes", 7.8, 6.58, 4.9, 0.36, theme, size=10, color_value=theme.muted)
    add_footer(slide, 1, total, theme)


def add_framework_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    add_text(slide, "Five transferable levers", 0.7, 1.18, 6.3, 0.5, theme, size=18, bold=True, color_value=theme.accent)
    card_w = 2.25
    for item_index, bullet in enumerate(spec.bullets):
        x = 0.7 + item_index * 2.48
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.72, 2.02, 0.78, 0.78, theme.accent if item_index % 2 == 0 else theme.accent2)
        add_text(slide, str(item_index + 1), x + 0.72, 2.16, 0.78, 0.28, theme, size=16, bold=True, color_value=color("#ffffff"), align=PP_ALIGN.CENTER)
        add_card(slide, x, 3.02, card_w, 1.55, theme, fill=theme.card, title=bullet, body=["Works across surfaces", "Tune for the task"], title_size=14, body_size=12)
    add_text(slide, "Use this map instead of teaching a separate framework for every Copilot surface.", 0.75, 5.55, 11.8, 0.45, theme, size=18, bold=True, color_value=theme.title)
    add_footer(slide, index, total, theme)


def add_surface_matrix_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    add_text(slide, "Not all controls exist everywhere. Teach the habit, then pick the right surface.", 0.7, 1.15, 10.8, 0.38, theme, size=16, color_value=theme.body)
    for item_index, bullet in enumerate(spec.bullets):
        name, _, detail = bullet.partition(":")
        row = item_index // 3
        col = item_index % 3
        x = 0.75 + col * 4.15
        y = 1.86 + row * 2.0
        add_card(slide, x, y, 3.7, 1.48, theme, fill=theme.card, title=name.strip(), body=[detail.strip()], title_size=16, body_size=13)
    add_card(slide, 4.9, 5.8, 3.7, 0.72, theme, fill=theme.card_alt, title="Facilitator cue", body=["Use CLI as the visible reference implementation."], title_size=13, body_size=11)
    add_footer(slide, index, total, theme)


def add_lever_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    lever_match = re.search(r"Lever (\d+)", spec.title)
    lever_number = lever_match.group(1) if lever_match else str(index)
    add_shape(slide, MSO_SHAPE.OVAL, 0.78, 1.54, 1.42, 1.42, theme.accent)
    add_text(slide, lever_number, 0.78, 1.86, 1.42, 0.42, theme, size=26, bold=True, color_value=color("#ffffff"), align=PP_ALIGN.CENTER, display=True)
    add_text(slide, "Token optimization lever", 2.45, 1.7, 4.1, 0.35, theme, size=13, bold=True, color_value=theme.accent)
    add_text(slide, spec.title.split(" - ", 1)[-1], 2.43, 2.02, 5.1, 0.62, theme, size=28, bold=True, color_value=theme.title, display=True)
    add_card(slide, 0.78, 3.25, 5.85, 2.55, theme, fill=theme.card, title="What students should practice", body=spec.bullets[:3], title_size=15, body_size=16)
    add_card(slide, 7.05, 1.7, 5.35, 4.1, theme, fill=theme.card_alt, title="Classroom move", body=spec.bullets[3:] or spec.bullets[:2], title_size=15, body_size=16)
    add_footer(slide, index, total, theme)


def guidance_for_pattern(title: str) -> list[str]:
    lower = title.lower()
    if "vs code" in lower:
        return ["Best for daily coding flow", "Use Ask, Plan, Agent, then Review deliberately"]
    if "github.com" in lower:
        return ["Best when page context matters", "Start from the issue, PR, discussion, or repo page"]
    if "cli" in lower:
        return ["Best for transparent tool control", "Filter noisy command output before it enters context"]
    if "coding agent" in lower:
        return ["Best for scoped async implementation", "Write the task like a small implementation brief"]
    if "code review" in lower:
        return ["Best for PR-shaped feedback", "Tune PR size and review instructions first"]
    return ["Pick the surface based on context boundaries", "Measure retries and rework"]


def add_pattern_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    add_card(slide, 0.78, 1.35, 7.0, 4.95, theme, fill=theme.card, title="Recommended pattern", body=spec.bullets, title_size=16, body_size=17)
    add_card(slide, 8.15, 1.35, 4.25, 2.1, theme, fill=theme.card_alt, title="Surface fit", body=guidance_for_pattern(spec.title), title_size=15, body_size=15)
    add_card(slide, 8.15, 3.75, 4.25, 2.05, theme, fill=theme.card, title="Risk to watch", body=["Vague scope expands context", "Broad context creates broad answers"], title_size=15, body_size=15)
    add_footer(slide, index, total, theme)


def add_exercise_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    add_text(slide, "Hands-on exercise", 0.75, 1.16, 3.0, 0.33, theme, size=13, bold=True, color_value=theme.accent)
    for item_index, bullet in enumerate(spec.bullets):
        x = 0.78 + (item_index % 2) * 5.95
        y = 1.85 + (item_index // 2) * 2.05
        add_shape(slide, MSO_SHAPE.OVAL, x, y, 0.55, 0.55, theme.accent2 if item_index % 2 else theme.accent)
        add_text(slide, str(item_index + 1), x, y + 0.11, 0.55, 0.22, theme, size=12, bold=True, color_value=color("#ffffff"), align=PP_ALIGN.CENTER)
        add_card(slide, x + 0.72, y - 0.02, 4.9, 1.18, theme, fill=theme.card, title=bullet, body=[], title_size=15)
    add_card(slide, 0.78, 6.0, 11.65, 0.48, theme, fill=theme.card_alt, title="Output", body=["A smaller, safer context plan students can reuse after the workshop."], title_size=12, body_size=11)
    add_footer(slide, index, total, theme)


def add_delivery_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    widths = [3.55, 3.55, 3.55]
    labels = ["Awareness", "Practitioner", "Applied review"]
    for item_index, bullet in enumerate(spec.bullets):
        x = 0.78 + item_index * 4.1
        add_card(slide, x, 1.75, widths[item_index], 3.35, theme, fill=theme.card, title=labels[item_index], body=[bullet], title_size=17, body_size=17)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 5.28, widths[item_index], 0.14, theme.accent if item_index != 2 else theme.accent2)
    add_text(slide, "Use preflight before the 4-hour version, especially when students review their own repositories.", 0.85, 6.0, 10.8, 0.35, theme, size=14, color_value=theme.body)
    add_footer(slide, index, total, theme)


def add_metric_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    for item_index, bullet in enumerate(spec.bullets):
        x = 0.82 + (item_index % 3) * 4.05
        y = 1.55 + (item_index // 3) * 1.8
        fill = theme.card_alt if item_index == 0 else theme.card
        add_card(slide, x, y, 3.55, 1.18, theme, fill=fill, title=bullet, body=[], title_size=16)
    add_card(slide, 2.4, 5.5, 8.3, 0.8, theme, fill=theme.card_alt, title="Promise measurement, not magic savings", body=["Use baselines, then compare retries, review quality, rework, and usage."], title_size=13, body_size=12)
    add_footer(slide, index, total, theme)


def add_takeaways_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 1.15, SLIDE_W, 0.16, theme.accent2)
    for item_index, bullet in enumerate(spec.bullets):
        y = 1.7 + item_index * 1.05
        add_card(slide, 1.05, y, 11.0, 0.74, theme, fill=theme.card if item_index % 2 else theme.card_alt, title=bullet, body=[], title_size=17)
    add_footer(slide, index, total, theme)


def add_default_slide(slide, spec: SlideSpec, theme: Theme, index: int, total: int) -> None:
    apply_background(slide, theme)
    add_header(slide, spec, theme)
    if len(spec.bullets) <= 3:
        add_card(slide, 0.8, 1.55, 7.4, 4.25, theme, fill=theme.card, title="Key message", body=spec.bullets, title_size=16, body_size=19)
        add_card(slide, 8.65, 1.55, 3.65, 4.25, theme, fill=theme.card_alt, title="Remember", body=["Better context improves quality, speed, and cost control."], title_size=15, body_size=16)
    else:
        for item_index, bullet in enumerate(spec.bullets):
            x = 0.82 + (item_index % 2) * 5.85
            y = 1.35 + (item_index // 2) * 1.25
            add_card(slide, x, y, 5.35, 0.88, theme, fill=theme.card, title=bullet, body=[], title_size=15)
    add_footer(slide, index, total, theme)


def render_slide(prs: Presentation, spec: SlideSpec, theme: Theme, index: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    lower_title = spec.title.lower()

    if index == 1:
        add_title_slide(slide, spec, theme, total)
    elif spec.title == "The five levers":
        add_framework_slide(slide, spec, theme, index, total)
    elif spec.title == "Surface matrix":
        add_surface_matrix_slide(slide, spec, theme, index, total)
    elif lower_title.startswith("lever "):
        add_lever_slide(slide, spec, theme, index, total)
    elif lower_title.endswith("pattern"):
        add_pattern_slide(slide, spec, theme, index, total)
    elif "exercise" in lower_title:
        add_exercise_slide(slide, spec, theme, index, total)
    elif spec.title == "Delivery tracks":
        add_delivery_slide(slide, spec, theme, index, total)
    elif spec.title == "What improvement looks like":
        add_metric_slide(slide, spec, theme, index, total)
    elif spec.title == "Takeaways":
        add_takeaways_slide(slide, spec, theme, index, total)
    else:
        add_default_slide(slide, spec, theme, index, total)

    add_notes(slide, spec)


def build_deck(theme: Theme, slides: list[SlideSpec]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "Token Optimization and Context Engineering"
    prs.core_properties.subject = theme.label
    prs.core_properties.author = "GitHub Copilot"
    prs.core_properties.created = datetime(2026, 1, 1, tzinfo=timezone.utc)
    prs.core_properties.modified = datetime(2026, 1, 1, tzinfo=timezone.utc)

    for index, spec in enumerate(slides, start=1):
        render_slide(prs, spec, theme, index, len(slides))

    output_path = DECKS_DIR / theme.output_name
    prs.save(output_path)
    validate_deck(output_path, slides)
    return output_path


def normalize(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def validate_deck(path: Path, slides: list[SlideSpec]) -> None:
    deck = Presentation(path)
    if len(deck.slides) != len(slides):
        raise ValueError(f"{path.name} has {len(deck.slides)} slides, expected {len(slides)}")

    for index, (slide, spec) in enumerate(zip(deck.slides, slides, strict=True), start=1):
        note_text = normalize(slide.notes_slide.notes_text_frame.text)
        expected = normalize(spec.notes)
        if expected not in note_text:
            raise ValueError(f"{path.name} slide {index} is missing expected speaker notes")

    with zipfile.ZipFile(path) as package:
        note_slide_parts = [
            name
            for name in package.namelist()
            if name.startswith("ppt/notesSlides/notesSlide") and name.endswith(".xml")
        ]
    if len(note_slide_parts) != len(slides):
        raise ValueError(f"{path.name} has {len(note_slide_parts)} notes slides, expected {len(slides)}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--variant",
        choices=[*THEMES.keys(), "all"],
        default="all",
        help="Deck variant to generate. Defaults to all variants.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    slides = load_slides()
    variants = THEMES.keys() if args.variant == "all" else [args.variant]

    for variant in variants:
        output = build_deck(THEMES[variant], slides)
        print(f"Generated {output.relative_to(REPO_ROOT)}")


if __name__ == "__main__":
    main()
